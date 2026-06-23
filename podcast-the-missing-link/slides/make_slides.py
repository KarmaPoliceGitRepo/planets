#!/usr/bin/env python3
"""make_slides.py — build the "How to use The Missing Link Studio" deck.

Produces an editable PowerPoint and (if LibreOffice is present) a matching PDF
that walk a complete beginner from raw video+audio all the way to a published
podcast on YouTube and Spotify.

    python3 slides/make_slides.py

Outputs (next to this script):
    slides/how-to-use-the-missing-link.pptx
    slides/how-to-use-the-missing-link.pdf   (if `soffice`/`libreoffice` exists)

Re-run any time you change the workflow; it overwrites the files.
"""
from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    print("Install python-pptx first:  python3 -m pip install python-pptx")
    raise SystemExit(2)

HERE = Path(__file__).resolve().parent
OUT_PPTX = HERE / "how-to-use-the-missing-link.pptx"
OUT_PDF = HERE / "how-to-use-the-missing-link.pdf"

# A plain-text outline of every slide. The rich .pptx is hand-built below; this
# same outline drives the reportlab PDF fallback so a PDF exists on any machine
# (even where LibreOffice can't convert the .pptx).
OUTLINE = [
    ("The Missing Link — Studio",
     ["Turn raw video + audio into a finished podcast — on your own computer, for free.",
      "You bring a camera and a microphone; everything else is set up for you.",
      "Free tools a 12-year-old can use: OBS, Audacity, FFmpeg, Whisper, Canva,",
      "Spotify for Creators, YouTube — plus the one-window Studio app.",
      "The show: why Nepalis who love their country still migrate — villages, energy, tourism."]),
    ("The journey of one episode",
     ["1. Record  — camera + mic into OBS; files land in raw/.",
      "2. Clean   — one click removes hiss, rumble, clicks, harsh s's.",
      "3. Choose  — auto-split into topics; tick what to keep.",
      "4. Master  — auto -16 LUFS loudness, MP3 + MP4, PASS/FAIL.",
      "5. Caption — Whisper writes the transcript + subtitles.",
      "6. Publish — Spotify for Creators + YouTube. Done.",
      "You never touch the command line — one window runs every step."]),
    ("Before you start (one time only)",
     ["1. Install 3 free apps: OBS Studio, Audacity, Python.",
      "2. Run the helper:  bash scripts/setup.sh  (checks FFmpeg, installs Whisper in a .venv).",
      "3. Make 3 free accounts with an adult: Spotify for Creators, YouTube, Drive/MEGA.",
      "4. Start the Studio:  python3 scripts/studio.py  → open http://127.0.0.1:8765"]),
    ("The Studio window, button by button",
     ["① Clean audio — high-pass, FFT denoise, de-click, de-ess, compression (light/medium/strong).",
      "② Split into segments — Whisper transcribes; audio grouped into tagged topics & sub-sections.",
      "Tick what to keep — untick whole segments or single sub-sections; see kept length live.",
      "③ Rebuild (keep ticked) — FFmpeg re-joins only kept parts; audio + video stay in sync.",
      "④ Master (-16 LUFS) — two-pass loudness; exports episode.mp3 + episode.mp4; PASS/FAIL.",
      "Subtitles — Whisper writes episode.txt + episode.srt."]),
    ("Segments vs. sub-sections — the editing superpower",
     ["You edit by MEANING, not by scrubbing a waveform.",
      "SEGMENT 'kings · villages · roads' 04:12-09:30  [keep]",
      "   - 'return to the villages movement' 04:12-05:40  [keep]",
      "   - 'a rambling tangent about lunch' 05:40-06:55  [drop]",
      "   - 'roads, schools, clinics promised' 06:55-09:30  [keep]",
      "SEGMENT 'geothermal · tato pani · tourism' 09:30-14:05  [keep]",
      "Keep a topic, drop one tangent inside it, then press Rebuild. Nothing is lost."]),
    ("Publish everywhere — for free",
     ["Audio → Spotify for Creators: upload episode.mp3; it auto-builds your RSS feed;",
      "   1-click to Spotify; submit the RSS to Apple Podcasts. RSS keeps you portable.",
      "Video → YouTube: upload episode.mp4; add captions/episode.srt; reuse chapters;",
      "   add the 1400x1400 cover as the thumbnail.",
      "Then back up the whole episode folder to the cloud — never keep only one copy."]),
    ("Cheat-sheet (terminal optional — all of these are also buttons)",
     ["Start the Studio:        python3 scripts/studio.py",
      'New episode:             bash scripts/new_episode.sh 01 "The Missing Link"',
      "Clean audio:             bash scripts/clean_audio.sh episodes/01-the-missing-link medium",
      "Split into segments:     python3 scripts/segment_episode.py episodes/01-the-missing-link",
      "Rebuild from choices:    python3 scripts/render_selection.py episodes/01-the-missing-link",
      "Master to -16 LUFS:      bash scripts/process_episode.sh episodes/01-the-missing-link",
      "Subtitles:               python3 scripts/transcribe.py episodes/01-the-missing-link"]),
    ("Done is better than perfect",
     ["Record · Clean · Choose · Master · Caption · Publish — then do it again next week.",
      "Full guides live in 02-implementation/ · the show's story lives in 03-content/."]),
]

# palette (matches the Studio UI)
BG = RGBColor(0x0F, 0x17, 0x20)
CARD = RGBColor(0x16, 0x21, 0x2E)
INK = RGBColor(0xE8, 0xEE, 0xF5)
MUT = RGBColor(0x90, 0xA4, 0xB8)
ACC = RGBColor(0x39, 0xB8, 0x8B)
ACC2 = RGBColor(0x2F, 0x7F, 0xD1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, left, top, width, height, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(left), Inches(top), Inches(width), Inches(height))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def text(s, left, top, width, height, runs, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, 18, INK, False)]
    first = True
    for item in runs:
        txt, size, color, bold = (item + (False,))[:4] if len(item) < 4 else item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
        r.font.name = "Calibri"
    return tb


def title_slide():
    s = slide()
    box(s, 0, 0, 13.333, 2.4, fill=RGBColor(0x12, 0x3C, 0x33))
    text(s, 0.9, 0.7, 11.5, 1.4,
         [("🎙️  The Missing Link — Studio", 40, INK, True),
          ("Turn raw video + audio into a finished podcast — on your own computer, for free.",
           20, MUT, False)])
    text(s, 0.9, 3.0, 11.5, 3.0,
         [("What this is", 22, ACC, True),
          ("A complete, beginner-friendly podcast workshop. You bring a camera and a "
           "microphone; everything else — recording, cleaning, editing-by-choosing, "
           "mastering to broadcast loudness, subtitles and publishing — is set up for you.",
           18, INK, False),
          ("", 8, INK, False),
          ("Built with free tools a 12-year-old can use: OBS, Audacity, FFmpeg, Whisper, "
           "Canva, Spotify for Creators, YouTube — plus the one-window Studio app.",
           16, MUT, False)])
    text(s, 0.9, 6.7, 11.5, 0.5, [("The show: why Nepalis who love their country still migrate — villages, energy & high-value tourism.", 14, MUT, False)])


def big_picture():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("The journey of one episode", 30, INK, True)])
    steps = [
        ("1", "Record", "Camera + mic into OBS.\nFiles land in raw/.", ACC2),
        ("2", "Clean", "One click removes hiss,\nrumble, clicks, harsh s’s.", ACC),
        ("3", "Choose", "Auto-split into topics;\ntick what to keep.", ACC2),
        ("4", "Master", "Auto -16 LUFS loudness,\nMP3 + MP4, PASS/FAIL.", ACC),
        ("5", "Caption", "Whisper writes the\ntranscript + subtitles.", ACC2),
        ("6", "Publish", "Spotify for Creators\n+ YouTube. Done.", ACC),
    ]
    x = 0.7
    w = 1.95
    for n, t, d, c in steps:
        box(s, x, 1.7, w, 3.2, fill=CARD, line=c)
        text(s, x, 1.9, w, 0.8, [(n, 30, c, True)], align=PP_ALIGN.CENTER)
        text(s, x, 2.7, w, 0.6, [(t, 18, INK, True)], align=PP_ALIGN.CENTER)
        text(s, x + 0.08, 3.4, w - 0.16, 1.4, [(d, 13, MUT, False)], align=PP_ALIGN.CENTER)
        x += w + 0.13
    box(s, 0.7, 5.4, 11.93, 1.4, fill=RGBColor(0x12, 0x2A, 0x22), line=ACC)
    text(s, 1.0, 5.6, 11.4, 1.1,
         [("The big idea: you never touch the command line. One window — the Studio — "
           "runs every step with buttons, and lets you keep exactly the bits you like.",
           18, INK, True)])


def one_time_setup():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("Before you start (one time only)", 30, INK, True)])
    rows = [
        ("Install 3 free apps", "OBS Studio (record), Audacity (optional fine edits), Python (runs the Studio)."),
        ("Run the setup helper", "In a terminal:  bash scripts/setup.sh  — it checks FFmpeg and installs Whisper into a tidy .venv."),
        ("Make 3 free accounts (with an adult)", "Spotify for Creators (host), YouTube (video), Google Drive/MEGA (backup)."),
        ("Start the Studio", "python3 scripts/studio.py  → open http://127.0.0.1:8765 in any browser."),
    ]
    y = 1.7
    for i, (h, d) in enumerate(rows, 1):
        box(s, 0.8, y, 11.7, 1.15, fill=CARD, line=None)
        text(s, 1.0, y + 0.12, 0.7, 0.9, [(str(i), 26, ACC, True)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.8, y + 0.12, 10.5, 0.95,
             [(h, 18, INK, True), (d, 15, MUT, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += 1.32


def studio_tour():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("The Studio window, button by button", 30, INK, True)])
    items = [
        ("① Clean audio", "FFmpeg DSP chain: high-pass 80 Hz, FFT denoise, de-click, de-ess, gentle compression. Pick light / medium / strong.", ACC),
        ("② Split into segments", "Whisper transcribes, then the audio is grouped into topic SEGMENTS and smaller SUB-SECTIONS, each auto-tagged.", ACC2),
        ("Tick what to keep", "Untick a whole segment, or single sub-sections inside it. The bar shows how long your kept episode will be.", INK),
        ("③ Rebuild (keep ticked)", "FFmpeg trims and re-joins only the parts you kept — audio and video stay perfectly in sync.", ACC),
        ("④ Master (-16 LUFS)", "Two-pass loudness normalisation to the podcast standard, exports episode.mp3 + episode.mp4, prints PASS/FAIL.", ACC2),
        ("Subtitles", "Whisper writes episode.txt + episode.srt for YouTube and accessibility.", INK),
    ]
    y = 1.6
    for h, d, c in items:
        box(s, 0.8, y, 11.7, 0.82, fill=CARD, line=None)
        text(s, 1.0, y + 0.06, 3.4, 0.7, [(h, 16, c, True)], anchor=MSO_ANCHOR.MIDDLE)
        text(s, 4.5, y + 0.06, 7.8, 0.7, [(d, 13, MUT, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.92


def segment_explainer():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("Segments vs. sub-sections — the editing superpower", 28, INK, True)])
    text(s, 0.9, 1.4, 11.5, 0.9,
         [("Instead of scrubbing a waveform, you edit by MEANING. The episode is cut into "
           "topics you can read and pick from.", 17, MUT, False)])
    box(s, 0.8, 2.4, 11.7, 3.7, fill=CARD, line=ACC2)
    text(s, 1.1, 2.6, 11.2, 3.4,
         [("📁  SEGMENT — “kings · villages · roads”   04:12–09:30      ☑ keep", 18, ACC, True),
          ("      ▫ sub-section  “return to the villages movement”   04:12–05:40   ☑", 15, INK, False),
          ("      ▫ sub-section  “a rambling tangent about lunch”     05:40–06:55   ☐  (untick!)", 15, MUT, False),
          ("      ▫ sub-section  “roads, schools, clinics promised”   06:55–09:30   ☑", 15, INK, False),
          ("", 8, INK, False),
          ("📁  SEGMENT — “geothermal · tato pani · tourism”   09:30–14:05    ☑ keep", 18, ACC, True),
          ("      ▫ sub-section  “tato pani springs heat homes”      09:30–11:20   ☑", 15, INK, False),
          ("      ▫ sub-section  “high-value tourism in Mustang”     11:20–14:05   ☑", 15, INK, False)])
    text(s, 0.9, 6.3, 11.5, 0.8,
         [("Keep a whole segment, drop one tangent inside it, or remove an entire topic — "
           "then press Rebuild. Nothing original is lost; re-choose any time.", 15, MUT, False)])


def publish_slide():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("Publish everywhere — for free", 30, INK, True)])
    cols = [
        ("🎧 Audio → Spotify for Creators", ACC, [
            "Upload exports/episode.mp3",
            "Paste title, notes, chapters",
            "It auto-builds your RSS feed",
            "1-click to Spotify; submit RSS to Apple",
            "RSS = portable: never lose your audience",
        ]),
        ("📺 Video → YouTube", ACC2, [
            "Upload exports/episode.mp4",
            "Add captions/episode.srt (subtitles)",
            "Reuse chapters as timestamps",
            "Add the 1400×1400 cover as thumbnail",
            "Huge discovery + accessibility",
        ]),
    ]
    x = 0.8
    for h, c, lines in cols:
        box(s, x, 1.7, 5.85, 4.6, fill=CARD, line=c)
        text(s, x + 0.3, 1.95, 5.3, 0.6, [(h, 20, c, True)])
        runs = [("•  " + ln, 16, INK, False) for ln in lines]
        text(s, x + 0.3, 2.8, 5.3, 3.2, runs)
        x += 6.05
    box(s, 0.8, 6.5, 11.85, 0.8, fill=RGBColor(0x12, 0x2A, 0x22), line=ACC)
    text(s, 1.0, 6.6, 11.4, 0.6,
         [("Then back up the whole episode folder to the cloud — never keep only one copy.", 16, INK, True)])


def cheatsheet():
    s = slide()
    text(s, 0.9, 0.5, 11.5, 0.8, [("Cheat-sheet (terminal optional)", 30, INK, True)])
    text(s, 0.9, 1.35, 11.5, 0.5, [("Everything below is also a button in the Studio window.", 15, MUT, False)])
    cmds = [
        ("Start the Studio (the easy way)", "python3 scripts/studio.py"),
        ("New episode folder", 'bash scripts/new_episode.sh 01 "The Missing Link"'),
        ("Clean the audio", "bash scripts/clean_audio.sh episodes/01-the-missing-link medium"),
        ("Split into tagged segments", "python3 scripts/segment_episode.py episodes/01-the-missing-link"),
        ("Rebuild from your choices", "python3 scripts/render_selection.py episodes/01-the-missing-link"),
        ("Master to -16 LUFS (MP3+MP4)", "bash scripts/process_episode.sh episodes/01-the-missing-link"),
        ("Subtitles", "python3 scripts/transcribe.py episodes/01-the-missing-link"),
    ]
    y = 1.95
    for h, c in cmds:
        text(s, 0.9, y, 5.2, 0.6, [(h, 15, INK, True)], anchor=MSO_ANCHOR.MIDDLE)
        box(s, 6.2, y, 6.4, 0.62, fill=RGBColor(0x0A, 0x11, 0x19), line=None)
        text(s, 6.4, y, 6.1, 0.62, [(c, 13, ACC, False)], anchor=MSO_ANCHOR.MIDDLE)
        y += 0.72


def closing():
    s = slide()
    box(s, 0, 2.6, 13.333, 2.3, fill=RGBColor(0x12, 0x3C, 0x33))
    text(s, 0.9, 3.0, 11.5, 1.6,
         [("Done is better than perfect.", 34, INK, True),
          ("Record · Clean · Choose · Master · Caption · Publish — then do it again next week.",
           20, MUT, False)],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 0.9, 5.3, 11.5, 0.6,
         [("Full guides live in 02-implementation/ · the show’s story lives in 03-content/.", 15, MUT, False)],
         align=PP_ALIGN.CENTER)


def pdf_via_libreoffice() -> bool:
    """Best-quality path: let LibreOffice render the .pptx to PDF 1:1."""
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return False
    print("[slides] trying LibreOffice for a pixel-perfect PDF…")
    try:
        rc = subprocess.run([soffice, "--headless", "--convert-to", "pdf",
                             "--outdir", str(HERE), str(OUT_PPTX)],
                            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                            timeout=120).returncode
    except Exception:
        return False
    return rc == 0 and OUT_PDF.exists()


def pdf_via_reportlab() -> bool:
    """Fallback: render the OUTLINE straight to a PDF (no LibreOffice needed)."""
    try:
        from reportlab.lib.pagesizes import landscape
        from reportlab.lib.units import inch
        from reportlab.pdfgen import canvas
    except ImportError:
        return False
    print("[slides] building PDF directly with reportlab…")
    W, H = landscape((13.333 * inch, 7.5 * inch))
    c = canvas.Canvas(str(OUT_PDF), pagesize=(W, H))

    def rgb(col):
        return (col[0] / 255, col[1] / 255, col[2] / 255)

    for i, (title, lines) in enumerate(OUTLINE):
        c.setFillColorRGB(*rgb(BG)); c.rect(0, 0, W, H, fill=1, stroke=0)
        accent = ACC if i % 2 == 0 else ACC2
        c.setFillColorRGB(*rgb(accent)); c.rect(0, H - 0.35 * inch, W, 0.35 * inch, fill=1, stroke=0)
        c.setFillColorRGB(*rgb(INK)); c.setFont("Helvetica-Bold", 26)
        c.drawString(0.8 * inch, H - 1.15 * inch, title)
        c.setFont("Helvetica", 15)
        y = H - 1.9 * inch
        for ln in lines:
            color = MUT if ln.strip().startswith("-") else INK
            c.setFillColorRGB(*rgb(color))
            font = "Courier" if ("scripts/" in ln or "bash " in ln) else "Helvetica"
            c.setFont(font, 14 if font == "Courier" else 15)
            c.drawString(1.0 * inch, y, ln[:120])
            y -= 0.62 * inch
        c.setFillColorRGB(*rgb(MUT)); c.setFont("Helvetica-Oblique", 10)
        c.drawRightString(W - 0.6 * inch, 0.45 * inch, f"The Missing Link — Studio   ·   {i + 1}/{len(OUTLINE)}")
        c.showPage()
    c.save()
    return OUT_PDF.exists()


def main() -> int:
    title_slide()
    big_picture()
    one_time_setup()
    studio_tour()
    segment_explainer()
    publish_slide()
    cheatsheet()
    closing()
    prs.save(str(OUT_PPTX))
    print(f"✅ Saved {OUT_PPTX}  ({len(prs.slides._sldIdLst)} slides)")

    if pdf_via_libreoffice():
        print(f"✅ Saved {OUT_PDF}  (LibreOffice, matches the .pptx exactly)")
    elif pdf_via_reportlab():
        print(f"✅ Saved {OUT_PDF}  (reportlab fallback — same content as the deck)")
    else:
        print("ℹ️  Could not make a PDF here. Open the .pptx and File → Export → PDF,")
        print("    or:  python3 -m pip install reportlab  then re-run this script.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
